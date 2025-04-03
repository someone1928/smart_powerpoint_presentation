import cv2
import mediapipe as mp
import pyautogui
import pyttsx3
import os
import time
import threading
from pptx import Presentation

ppt_file = "presentation.pptx"
if not os.path.exists(ppt_file):
    print("Error: PowerPoint file not found!")
    exit()

os.startfile(ppt_file)

mp_hands = mp.solutions.hands
mp_drawing = mp.solutions.drawing_utils
hands = mp_hands.Hands(min_detection_confidence=0.7, min_tracking_confidence=0.7)

engine = pyttsx3.init()
engine.setProperty("rate", 150)

presentation = Presentation(ppt_file)
current_slide = 0

cap = cv2.VideoCapture(0)
last_action_time = 0
reading_thread = None

def read_slide(slide_number):
    global reading_thread
    if 0 <= slide_number < len(presentation.slides):
        slide = presentation.slides[slide_number]
        text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
        if text:
            if reading_thread is None or not reading_thread.is_alive():
                reading_thread = threading.Thread(target=lambda: (engine.say(text), engine.runAndWait()))
                reading_thread.start()

while cap.isOpened():
    ret, frame = cap.read()
    if not ret:
        break

    frame = cv2.flip(frame, 1)
    rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
    results = hands.process(rgb_frame)

    if results.multi_hand_landmarks:
        for hand_landmarks in results.multi_hand_landmarks:
            mp_drawing.draw_landmarks(frame, hand_landmarks, mp_hands.HAND_CONNECTIONS)

            landmarks = hand_landmarks.landmark
            index_tip = landmarks[8]
            thumb_tip = landmarks[4]
            middle_tip = landmarks[12]
            ring_tip = landmarks[16]
            pinky_tip = landmarks[20]
            wrist = landmarks[0]

            current_time = time.time()

            if index_tip.y < middle_tip.y and current_time - last_action_time > 0.8:
                print("Next Slide")
                pyautogui.press("right")
                current_slide = min(len(presentation.slides) - 1, current_slide + 1)
                last_action_time = current_time

            elif thumb_tip.y < wrist.y and current_time - last_action_time > 0.8:
                print("Previous Slide")
                pyautogui.press("left")
                current_slide = max(0, current_slide - 1)
                last_action_time = current_time

            if abs(index_tip.y - middle_tip.y) < 0.02:
                print("Reading Slide")
                read_slide(current_slide)
                last_action_time = current_time

            if index_tip.y < middle_tip.y and middle_tip.y < ring_tip.y:
                print("Start Slideshow")
                pyautogui.press("f5")
                last_action_time = current_time

            if (abs(index_tip.x - pinky_tip.x) > 0.1) and (abs(index_tip.y - pinky_tip.y) < 0.05):
                print("Pause Presentation")
                pyautogui.press("space")
                last_action_time = current_time

            distance_thumb_index = abs(thumb_tip.x - index_tip.x) + abs(thumb_tip.y - index_tip.y)
            if distance_thumb_index < 0.05:
                print("Zooming In")
                pyautogui.keyDown("ctrl")
                pyautogui.scroll(2)
                pyautogui.keyUp("ctrl")
                last_action_time = current_time
            elif distance_thumb_index > 0.2:
                print("Zooming Out")
                pyautogui.keyDown("ctrl")
                pyautogui.scroll(-2)
                pyautogui.keyUp("ctrl")
                last_action_time = current_time

    cv2.imshow("Hand Gesture PPT Control", frame)

    if cv2.waitKey(1) & 0xFF == ord("q"):
        break

cap.release()
cv2.destroyAllWindows()
